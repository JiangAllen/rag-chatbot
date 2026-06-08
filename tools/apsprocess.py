from tools.datapreprocessing import datapreprocessing
from pathlib import Path
from pptx import Presentation
import json
import re
import os

def classify_node_type(text):
    text = text.strip()
    if not text:
        return "unknown"
    if re.search(r'開始|Start', text):
        return "start"
    if re.search(r'結束|End', text):
        return "end"
    if re.search(r'是否|？|\?', text) or re.match(r'^[YN]$', text):
        return "decision"
    return "process"

def parse_slide_nodes(slide, slide_num):
    nodes = []
    node_id_counter = 1
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text.strip()
        if not text:
            continue
        node_type = classify_node_type(text)
        node = {
            "id": f"slide{slide_num}_node{node_id_counter}",
            "text": text,
            "type": node_type,
            "source_slide": slide_num,
            "position": {
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height
            }
        }
        nodes.append(node)
        node_id_counter += 1
    return nodes

def infer_edges(nodes):
    sorted_nodes = sorted(nodes, key=lambda n: (n["position"]["top"], n["position"]["left"]))
    edges = []
    for i in range(len(sorted_nodes) - 1):
        edges.append({
            "from": sorted_nodes[i]["id"],
            "to": sorted_nodes[i + 1]["id"]
        })
    return edges

def parse_pptx_by_slide(pptx_path, output_dir):
    prs = Presentation(pptx_path)
    os.makedirs(output_dir, exist_ok=True)
    for slide_num, slide in enumerate(prs.slides, start=1):
        nodes = parse_slide_nodes(slide, slide_num)
        edges = infer_edges(nodes)
        flowchart = {
            "title": f"Slide {slide_num} 流程圖",
            "nodes": nodes,
            "edges": edges
        }
        outpath = os.path.join(output_dir, f"slide_{slide_num}.json")
        with open(outpath, "w", encoding="utf-8") as f:
            json.dump(flowchart, f, ensure_ascii=False, indent=2)
        print(f"✅ 已輸出：{outpath}")

def convert_folder_to_images(folder_path):
    folder = Path(folder_path)
    for file in folder.glob("*"):
        if not file.is_file():
            continue
        output_dir = folder / "{}_images".format(file.stem)
        output_dir.mkdir(exist_ok=True)
        if file.suffix.lower() == ".pdf":
            dt.convert_pdf_image(file, output_dir)

if __name__ == "__main__":
    dt = datapreprocessing()
    pptx_path = "./data/ppts/福耀.pptx"
    output_dir = "./data/ppts/json"
    parse_pptx_by_slide(pptx_path, output_dir)
    # convert_folder_to_images(pdf_path)
    result = dt.excel_to_json(file_path = "D:\Python\ZERONE\data\excels\Wrangling Input.xlsx")
    print(result)